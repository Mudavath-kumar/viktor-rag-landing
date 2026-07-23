"""
RAG pipeline using:
- sentence-transformers BAAI/bge-m3 for embeddings (1024-dim)
- Qdrant Cloud for vector storage and search
- xAI Grok API for text generation (OpenAI-compatible endpoint)
"""
import json
import traceback
from typing import List, Optional, Any

import numpy as np
from openai import OpenAI
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance, VectorParams, PointStruct,
    Filter, FieldCondition, MatchValue, SearchRequest
)

import config
import data_store as store

# ─── Lazy singletons ──────────────────────────────────────────────────────────

_embed_model: Any = None
_qdrant: Optional[QdrantClient] = None
_xai: Optional[OpenAI] = None


def get_embed_model():
    global _embed_model
    if _embed_model is None:
        print(f"Loading embedding model: {config.EMBED_MODEL}...")
        from sentence_transformers import SentenceTransformer
        # Force device='cpu' to avoid meta tensor errors with lazy loading
        _embed_model = SentenceTransformer(
            config.EMBED_MODEL,
            device="cpu",
            model_kwargs={"device_map": None},
        )
        print(f"Embedding model '{config.EMBED_MODEL}' ready (dim={_embed_model.get_sentence_embedding_dimension()}).")
    return _embed_model


def get_qdrant() -> QdrantClient:
    global _qdrant
    if _qdrant is None:
        _qdrant = QdrantClient(
            url=config.QDRANT_HOST,
            api_key=config.QDRANT_API_KEY,
            timeout=30,
        )
        _ensure_collection(_qdrant)
    return _qdrant


def get_xai() -> OpenAI:
    global _xai
    if _xai is None:
        if not config.XAI_API_KEY:
            raise ValueError("XAI_API_KEY is not set")
        _xai = OpenAI(
            api_key=config.XAI_API_KEY,
            base_url="https://api.x.ai/v1",
        )
    return _xai


def _ensure_collection(client: QdrantClient):
    """Create Qdrant collection if it does not exist, or recreate if dimensions mismatch."""
    from qdrant_client.models import PayloadSchemaType
    existing = [c.name for c in client.get_collections().collections]
    if config.QDRANT_COLLECTION in existing:
        # Check vector dimensions match current model
        info = client.get_collection(config.QDRANT_COLLECTION)
        actual_size = info.config.params.vectors.size
        if actual_size != config.VECTOR_SIZE:
            print(f"Vector dimension mismatch: collection has {actual_size}-dim, config wants {config.VECTOR_SIZE}-dim. Recreating...")
            client.delete_collection(config.QDRANT_COLLECTION)
        else:
            # Ensure payload indexes exist even if collection already exists
            try:
                client.create_payload_index(collection_name=config.QDRANT_COLLECTION, field_name="user_id", field_schema=PayloadSchemaType.KEYWORD)
                client.create_payload_index(collection_name=config.QDRANT_COLLECTION, field_name="doc_id", field_schema=PayloadSchemaType.KEYWORD)
            except Exception as e:
                pass  # Ignore if indexes already exist
            return

    client.create_collection(
        collection_name=config.QDRANT_COLLECTION,
        vectors_config=VectorParams(
            size=config.VECTOR_SIZE,
            distance=Distance.COSINE,
        ),
    )
    print(f"Created Qdrant collection: {config.QDRANT_COLLECTION} ({config.VECTOR_SIZE}-dim)")
    try:
        client.create_payload_index(collection_name=config.QDRANT_COLLECTION, field_name="user_id", field_schema=PayloadSchemaType.KEYWORD)
        client.create_payload_index(collection_name=config.QDRANT_COLLECTION, field_name="doc_id", field_schema=PayloadSchemaType.KEYWORD)
    except Exception:
        pass


# ─── Text extraction ──────────────────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""
    try:
        if ext == "txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        if ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            return "\n".join(
                (page.extract_text() or "") for page in reader.pages
            )
        if ext == "docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if ext in ("xlsx", "xls", "csv"):
            import pandas as pd
            if ext == "csv":
                df = pd.read_csv(file_path, encoding="utf-8", errors="ignore")
            else:
                df = pd.read_excel(file_path)
            return df.to_string(index=False)
        if ext == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        # Fallback: try reading as plain text
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(100_000)
    except Exception as e:
        print(f"extract_text error ({file_path}): {e}")
        return ""


# ─── Chunking ─────────────────────────────────────────────────────────────────

def chunk_text(text: str, max_chars: int = 800) -> List[str]:
    chunks, paragraphs = [], text.split("\n\n")
    current = ""
    for p in paragraphs:
        if len(current) + len(p) > max_chars and current:
            chunks.append(current.strip())
            current = p
        else:
            current = (current + "\n\n" + p) if current else p
    if current.strip():
        chunks.append(current.strip())
    return chunks or [text[:max_chars]]


# ─── Document processing ──────────────────────────────────────────────────────

async def process_document(doc_id: str, text: str, user_id: str):
    """Embed chunks and upsert to Qdrant asynchronously."""
    import asyncio
    model = get_embed_model()
    qdrant = get_qdrant()

    chunks = chunk_text(text)
    if not chunks:
        await store.update_document_status(user_id, doc_id, "error")
        return

    print(f"Embedding {len(chunks)} chunks for doc {doc_id}...")
    embeddings = await asyncio.to_thread(
        model.encode, chunks, normalize_embeddings=True, convert_to_numpy=True, show_progress_bar=False
    )

    points = [
        PointStruct(
            id=_chunk_id(doc_id, i),
            vector=embeddings[i].tolist(),
            payload={
                "doc_id": doc_id,
                "user_id": user_id,
                "chunk_index": i,
                "content": chunk,
            },
        )
        for i, chunk in enumerate(chunks)
    ]

    await asyncio.to_thread(qdrant.upsert, collection_name=config.QDRANT_COLLECTION, points=points, wait=True)
    await store.update_document_status(user_id, doc_id, "done")
    print(f"[OK] Doc {doc_id}: {len(chunks)} chunks indexed to Qdrant")


def _chunk_id(doc_id: str, chunk_index: int) -> str:
    """Generate a deterministic UUID-like integer id for Qdrant point."""
    import hashlib
    h = hashlib.md5(f"{doc_id}:{chunk_index}".encode()).hexdigest()
    return int(h, 16) % (2**63)


# ─── Search ───────────────────────────────────────────────────────────────────

def search_chunks(
    user_id: str, query: str, document_id: Optional[str] = None, limit: int = 5
) -> List[dict]:
    model = get_embed_model()
    qdrant = get_qdrant()

    query_emb = model.encode([query], normalize_embeddings=True, convert_to_numpy=True)[0]

    # Build filter: always by user, optionally by doc
    must_conditions = [
        FieldCondition(key="user_id", match=MatchValue(value=user_id))
    ]
    if document_id:
        must_conditions.append(
            FieldCondition(key="doc_id", match=MatchValue(value=document_id))
        )

    results = qdrant.query_points(
        collection_name=config.QDRANT_COLLECTION,
        query=query_emb.tolist(),
        query_filter=Filter(must=must_conditions),
        limit=limit,
        with_payload=True,
    ).points

    chunks = []
    for r in results:
        payload = r.payload or {}
        chunks.append({
            "content": payload.get("content", ""),
            "score": r.score,
            "doc_id": payload.get("doc_id", ""),
            "doc_name": payload.get("doc_name", "Unknown"),
        })
    return chunks


# ─── Enrich chunk with doc name ───────────────────────────────────────────────

async def _enrich_chunks_with_doc_name(chunks: list, user_id: str) -> list:
    """Add doc_name to chunks by looking up Document records."""
    import asyncio
    doc_cache = {}
    for c in chunks:
        doc_id = c.get("doc_id", "")
        if doc_id and doc_id not in doc_cache:
            d = await store.get_document(doc_id)
            doc_cache[doc_id] = d.get("name", "Unknown") if d else "Unknown"
        c["doc_name"] = doc_cache.get(doc_id, "Unknown")
    return chunks


# ─── Answer generation ────────────────────────────────────────────────────────

async def generate_answer(user_id: str, query: str, session_id: str) -> str:
    xai = get_xai()

    # Find session document scope
    document_id = None
    try:
        sessions = await store.get_sessions(user_id)
        for s in sessions:
            if s["id"] == session_id:
                document_id = s.get("document_id")
                break
    except Exception as e:
        print(f"Session lookup error: {e}")

    # Vector search
    raw_chunks = search_chunks(user_id, query, document_id=document_id, limit=5)
    chunks = await _enrich_chunks_with_doc_name(raw_chunks, user_id)

    if not chunks:
        return "No relevant documents found. Please upload some documents first."

    context = "\n\n---\n\n".join(
        f"[Source {i+1} — {c['doc_name']}]\n{c['content']}"
        for i, c in enumerate(chunks)
    )

    completion = xai.chat.completions.create(
        model=config.XAI_MODEL,
        temperature=config.XAI_TEMPERATURE,
        max_tokens=config.XAI_MAX_TOKENS,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a precise RAG assistant. Answer comprehensively based on the provided context. "
                    "Always cite which source the information came from. "
                    "If the context lacks information to answer, say so clearly."
                ),
            },
            {"role": "user", "content": f"Context:\n\n{context}\n\nQuestion: {query}"},
        ],
    )

    answer = completion.choices[0].message.content or "Sorry, could not generate an answer."

    # Append unique source references
    unique_docs = []
    for c in chunks:
        name = c.get("doc_name", "Unknown")
        if name not in unique_docs:
            unique_docs.append(name)

    if unique_docs:
        answer += "\n\n**Sources:** " + ", ".join(f"`{d}`" for d in unique_docs)

    return answer


# ─── Summary ──────────────────────────────────────────────────────────────────

def _get_doc_text_from_qdrant(doc_id: str, user_id: str, max_chunks: int = 15) -> str:
    qdrant = get_qdrant()
    # Scroll through stored chunks
    results, _ = qdrant.scroll(
        collection_name=config.QDRANT_COLLECTION,
        scroll_filter=Filter(must=[
            FieldCondition(key="doc_id", match=MatchValue(value=doc_id)),
            FieldCondition(key="user_id", match=MatchValue(value=user_id)),
        ]),
        limit=max_chunks,
        with_payload=True,
    )
    chunks = sorted(results, key=lambda p: p.payload.get("chunk_index", 0))
    return "\n\n".join(p.payload.get("content", "") for p in chunks)


async def generate_summary(doc_id: str, user_id: str) -> dict:
    xai = get_xai()
    text = _get_doc_text_from_qdrant(doc_id, user_id, max_chunks=15)
    if not text:
        return {"summary": "No content found.", "key_points": [], "topics": [], "tldr": "No content."}

    completion = xai.chat.completions.create(
        model=config.XAI_MODEL,
        temperature=0.1,
        max_tokens=800,
        messages=[
            {
                "role": "system",
                "content": (
                    'You are a document analysis AI. Return a JSON object with:\n'
                    '- "summary": 3-5 sentence overview\n'
                    '- "key_points": array of 4-6 key takeaways\n'
                    '- "topics": array of 3-5 main topics\n'
                    '- "tldr": single sentence TL;DR\n'
                    'Return ONLY valid JSON, no markdown, no code fences.'
                ),
            },
            {"role": "user", "content": f"Document content:\n\n{text}"},
        ],
    )

    raw = (completion.choices[0].message.content or "{}").strip()
    raw = _strip_json_fences(raw)
    try:
        result = json.loads(raw)
    except Exception:
        result = {"summary": raw, "key_points": [], "topics": [], "tldr": raw[:200]}

    summary_data = {
        "doc_id": doc_id,
        "summary": result.get("summary", ""),
        "key_points": result.get("key_points", []),
        "topics": result.get("topics", []),
        "tldr": result.get("tldr", ""),
    }
    await store.save_summary(doc_id, summary_data)

    # Auto-tag as well
    try:
        tags_result = await generate_tags(doc_id, user_id)
        summary_data["tags"] = tags_result.get("tags", [])
        summary_data["category"] = tags_result.get("category", "")
    except Exception as e:
        print(f"Auto-tag error: {e}")

    return summary_data


# ─── Quiz ─────────────────────────────────────────────────────────────────────

async def generate_quiz(doc_id: str, user_id: str) -> dict:
    xai = get_xai()
    text = _get_doc_text_from_qdrant(doc_id, user_id, max_chunks=12)
    if not text:
        return {"doc_id": doc_id, "questions": []}

    completion = xai.chat.completions.create(
        model=config.XAI_MODEL,
        temperature=0.3,
        max_tokens=1500,
        messages=[
            {
                "role": "system",
                "content": (
                    'You are a quiz generator AI. Create exactly 5 multiple-choice questions.\n'
                    'Return a JSON object with key "questions", each item having:\n'
                    '- "question": question text\n'
                    '- "options": array of exactly 4 answer strings\n'
                    '- "correct": index (0-3) of the correct answer\n'
                    '- "explanation": brief explanation\n'
                    'Return ONLY valid JSON, no markdown, no code fences.'
                ),
            },
            {"role": "user", "content": f"Create quiz from:\n\n{text}"},
        ],
    )

    raw = _strip_json_fences((completion.choices[0].message.content or "{}").strip())
    try:
        result = json.loads(raw)
    except Exception:
        result = {"questions": []}

    quiz_data = {"doc_id": doc_id, "questions": result.get("questions", [])}
    await store.save_quiz(doc_id, quiz_data)
    return quiz_data


# ─── Tags ─────────────────────────────────────────────────────────────────────

async def generate_tags(doc_id: str, user_id: str) -> dict:
    xai = get_xai()
    text = _get_doc_text_from_qdrant(doc_id, user_id, max_chunks=8)
    if not text:
        return {"doc_id": doc_id, "tags": [], "category": "Uncategorized"}

    completion = xai.chat.completions.create(
        model=config.XAI_MODEL,
        temperature=0.1,
        max_tokens=300,
        messages=[
            {
                "role": "system",
                "content": (
                    'You are a document classifier AI. Return a JSON object with:\n'
                    '- "tags": array of 3-6 descriptive tags\n'
                    '- "category": single primary category\n'
                    'Return ONLY valid JSON, no markdown, no code fences.'
                ),
            },
            {"role": "user", "content": f"Classify this document:\n\n{text[:3000]}"},
        ],
    )

    raw = _strip_json_fences((completion.choices[0].message.content or "{}").strip())
    try:
        result = json.loads(raw)
    except Exception:
        result = {"tags": ["Document"], "category": "General"}

    tags = result.get("tags", ["Document"])
    category = result.get("category", "General")
    await store.update_document_tags(user_id, doc_id, tags, category)
    return {"doc_id": doc_id, "tags": tags, "category": category}


# ─── Insights ─────────────────────────────────────────────────────────────────

async def generate_insights(user_id: str) -> dict:
    xai = get_xai()
    docs = await store.get_documents(user_id)
    if not docs:
        return {
            "key_topics": [], "knowledge_gaps": [], "connections": [],
            "briefing": "Upload some documents to generate insights.",
        }

    doc_summaries = []
    for d in docs[:20]:
        summary = await store.get_summary(d["id"])
        if summary and summary.get("tldr"):
            doc_summaries.append(f"Document '{d['name']}': {summary['tldr']}")
        else:
            text = _get_doc_text_from_qdrant(d["id"], user_id, max_chunks=2)
            if text:
                doc_summaries.append(f"Document '{d['name']}': {text[:200]}")

    if not doc_summaries:
        return {
            "key_topics": [], "knowledge_gaps": [], "connections": [],
            "briefing": "No document content available for analysis.",
        }

    completion = xai.chat.completions.create(
        model=config.XAI_MODEL,
        temperature=0.1,
        max_tokens=1000,
        messages=[
            {
                "role": "system",
                "content": (
                    'You are a knowledge base analyst. Given summaries of multiple documents, produce a JSON object with:\n'
                    '- "key_topics": array of 5-8 main topics across all documents\n'
                    '- "knowledge_gaps": array of 2-4 areas mentioned but not deeply covered\n'
                    '- "connections": array of 2-4 objects, each with "docs" (array of doc names) and "theme" (shared topic)\n'
                    '- "briefing": 3-4 sentence summary of the entire knowledge base\n'
                    'Return ONLY valid JSON, no markdown, no code fences.'
                ),
            },
            {"role": "user", "content": f"Analyze this knowledge base:\n\n{chr(10).join(doc_summaries)}"},
        ],
    )

    raw = _strip_json_fences((completion.choices[0].message.content or "{}").strip())
    try:
        result = json.loads(raw)
    except Exception:
        result = {"key_topics": [], "knowledge_gaps": [], "connections": [], "briefing": raw[:500]}

    insights = {
        "key_topics": result.get("key_topics", []),
        "knowledge_gaps": result.get("knowledge_gaps", []),
        "connections": result.get("connections", []),
        "briefing": result.get("briefing", ""),
    }
    await store.save_insights(user_id, insights)
    return insights


# ─── Util ─────────────────────────────────────────────────────────────────────

def _strip_json_fences(raw: str) -> str:
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[-1]
    if raw.endswith("```"):
        raw = raw.rsplit("```", 1)[0]
    return raw.strip()
